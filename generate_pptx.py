import sys
import subprocess

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation()

slides_data = [
    {
        "title": "CattleActiv-ML Intellektual Platformasi",
        "content": "Mavzu: Qoramollar sog'lig'ini va faolligini monitoring qilish uchun \"CattleActiv-ML\" intellektual platformasini ishlab chiqish.\n\nBajardi: Daribaev Adilbek.\n\nMaqsad: Chorvachilikda qoramollar sog'lig'ini real vaqt rejimida kuzatish, kasalliklarni barvaqt aniqlash va sun'iy intellekt yordamida tashxis qo'yish tizimini yaratish."
    },
    {
        "title": "Muammo va Dolzarblik",
        "content": "Muammo: Chorvachilikda kasalliklarni (masalan, Mastit, Ketoz) kech aniqlash oqibatida sut sog'imi kamayishi va iqtisodiy zararlar.\n\nYechim: IoT (Narsalar interneti) datchiklari va Machine Learning (Sun'iy intellekt) algoritmlarini birlashtirgan, fermerlar uchun qulay \"Gibrid AI\" platformasini joriy etish."
    },
    {
        "title": "Texnologik Stek",
        "content": "Backend: Python va Flask freymvorki.\n\nMa'lumotlar bazasi: SQLite (qoramollar profili, sog'im ko'rsatkichlari va kasalliklar tarixi uchun).\n\nFrontend: HTML, CSS (tabiatga xos zamonaviy \"Nature-inspired\" dizayn), JavaScript.\n\nAsosiy modullar: AI Diagnostika, IoT Simulyator, Ko'p tilli qobiq."
    },
    {
        "title": "Platformaning Asosiy Imkoniyatlari",
        "content": "1. Qoramollar Bazasini Boshqarish: Har bir sigir uchun individual profil, rasmlar, ID raqamlar va sut sog'ish unumdorligini kiritish/kuzatish.\n2. Real vaqtda monitoring: Har bir sigirning harorati, yurish faolligi va boshqa ko'rsatkichlarini kuzatib borish.\n3. AI Veterinar Yordamchi: Fermerlarga interaktiv maslahatlar beruvchi chat-bot tizimi.\n4. Bildirishnomalar (Notifications): Xavfli holatlar aniqlanganda tezkor ogohlantirishlar tizimi."
    },
    {
        "title": "Sun'iy Intellekt va Diagnostika Moduli",
        "content": "10 xil kasallikni aniqlash: Tizim simptomlar asosida 10 xil keng tarqalgan kasalliklarni interaktiv tarzda tahlil qiladi.\n\nQanday ishlaydi: Fermer simptomni kiritadi yoxud menyudan tanlaydi, AI esa kasallik nomi, asoratlari va davolash bo'yicha tibbiy ko'rsatmalarni dinamik tarzda chiqarib beradi."
    },
    {
        "title": "IoT Simulyatsiyasi (Innovatsion yondashuv)",
        "content": "Maqsad: Haqiqiy datchiklar o'rniga, tizimni sinovdan o'tkazish uchun aqlli simulyator yaratildi.\n\nQanday ishlaydi: Qoidalar asosida ishlovchi simulyatsiya (Rule-based IoT simulation) orqali sog'lom va kasal qoramollarning ma'lumotlarini haqiqiyga o'xshatib generatsiya qiladi va dashboard'ga uzatadi."
    },
    {
        "title": "Ko'p tilli qo'llab-quvvatlash va Dizayn",
        "content": "Lokalizatsiya: Tizim to'liq 4 ta tilda ishlaydi: O'zbek, Qoraqalpoq, Rus va Ingliz tillari.\n\nDizayn yechimlari: Fermerlar uchun tushunarli bo'lishi uchun zamonaviy animatsiyalar, qulay navigatsiya menyusi va ekologik toza (yashil) ranglar palitrasi ishlatildi."
    },
    {
        "title": "Xulosa va Kutilayotgan Natijalar",
        "content": "Platforma chorva mollarining kasallanish darajasini kamaytirishga va veterinarlarning ishini yengillashtirishga xizmat qiladi.\n\nOlingan natijalar shuni ko'rsatadiki, CattleActiv-ML orqali qilinadigan barvaqt diagnostika iqtisodiy samaradorlikni sezilarli darajada oshirishi mumkin."
    }
]

for slide_info in slides_data:
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 255, 240)
    
    title_shape = slide.shapes.title
    body_shape = slide.shapes.placeholders[1]
    
    title_shape.text = slide_info["title"]
    
    for paragraph in title_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.bold = True
            run.font.color.rgb = RGBColor(34, 139, 34)
            
    tf = body_shape.text_frame
    tf.text = slide_info["content"]
    
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(24)
        for run in paragraph.runs:
            run.font.color.rgb = RGBColor(50, 50, 50)

prs.save(r'd:\Diplom1\CattleActiv_Presentation.pptx')
print("Presentation generated successfully at d:\Diplom1\CattleActiv_Presentation.pptx")
